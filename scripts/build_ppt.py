"""
Build PersonalOS Agent hackathon PPT using python-pptx.
Run: python scripts/build_ppt.py
Output: docs/PersonalOS_Agent_Presentation.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Colour palette ─────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x06, 0x09, 0x12)   # deep navy
C_BG2       = RGBColor(0x0D, 0x14, 0x24)   # card bg
C_PURPLE    = RGBColor(0x8B, 0x5C, 0xF6)   # primary accent
C_PURPLE_HI = RGBColor(0xA7, 0x8B, 0xFA)   # lighter purple
C_INDIGO    = RGBColor(0x4F, 0x46, 0xE5)   # indigo
C_GREEN     = RGBColor(0x10, 0xB9, 0x81)   # green
C_ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)   # orange
C_RED       = RGBColor(0xEF, 0x44, 0x44)   # red
C_BLUE      = RGBColor(0x3B, 0x82, 0xF6)   # blue
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)   # white
C_TEXT      = RGBColor(0xE2, 0xE8, 0xF0)   # primary text
C_TEXT2     = RGBColor(0x94, 0xA3, 0xB8)   # secondary text
C_TEXT3     = RGBColor(0x47, 0x55, 0x69)   # muted text
C_BORDER    = RGBColor(0x1A, 0x25, 0x40)   # border

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ────────────────────────────────────────────────────────────────────

def blank_slide():
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()   # no border
    if radius:
        # round corners via XML
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, color=C_TEXT,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_para(tf, text, font_size=14, bold=False, color=C_TEXT2,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def accent_bar(slide, color=C_PURPLE, w=13.33):
    """Top 3px accent bar."""
    add_rect(slide, 0, 0, w, 0.042, color)


def gradient_title_bar(slide, color1=C_PURPLE, color2=C_INDIGO, h=1.15):
    """Gradient-like header using two overlapping rects."""
    add_rect(slide, 0, 0, 7.5, h, color1)
    add_rect(slide, 6.0, 0, 7.33, h, color2)


def slide_number_tag(slide, num, total=14):
    add_textbox(slide, f"{num} / {total}",
                12.5, 7.1, 0.8, 0.35,
                font_size=9, color=C_TEXT3, align=PP_ALIGN.RIGHT)


def tag_pill(slide, text, l, t, color=C_PURPLE, bg_alpha=None):
    """Coloured pill tag."""
    w = len(text) * 0.085 + 0.22
    r, g, b = color[0], color[1], color[2]
    bg_color = RGBColor(min(r + 20, 255), min(g + 15, 255), min(b + 40, 255))
    add_rect(slide, l, t, w, 0.28, bg_color, radius=40000)
    add_textbox(slide, text, l + 0.06, t + 0.03, w, 0.25,
                font_size=9, bold=True, color=color, align=PP_ALIGN.LEFT)
    return w


def card(slide, l, t, w, h, title, body_lines, title_color=C_PURPLE_HI, accent=None):
    """Card with optional left accent."""
    add_rect(slide, l, t, w, h, C_BG2, radius=20000)
    if accent:
        add_rect(slide, l, t, 0.045, h, accent)
    xl = l + (0.18 if accent else 0.14)
    add_textbox(slide, title, xl, t + 0.1, w - 0.3, 0.3,
                font_size=11, bold=True, color=title_color)
    y = t + 0.42
    for line in body_lines:
        add_textbox(slide, line, xl, y, w - 0.28, 0.28,
                    font_size=9.5, color=C_TEXT2)
        y += 0.26


def stat_box(slide, l, t, num, label, color=C_PURPLE_HI):
    add_rect(slide, l, t, 1.55, 1.1, C_BG2, radius=20000)
    add_textbox(slide, num,  l, t + 0.1, 1.55, 0.55,
                font_size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, l, t + 0.68, 1.55, 0.35,
                font_size=8.5, bold=True, color=C_TEXT3, align=PP_ALIGN.CENTER)


def flow_step(slide, l, t, text, color=C_PURPLE, width=3.8):
    add_rect(slide, l, t, width, 0.36, C_BG2, radius=20000)
    add_rect(slide, l, t + 0.09, 0.04, 0.18, color)
    add_textbox(slide, text, l + 0.14, t + 0.06, width - 0.2, 0.28,
                font_size=10, color=C_TEXT)


def arrow_down(slide, l, t):
    add_textbox(slide, "↓", l, t, 0.3, 0.3, font_size=14, color=C_TEXT3, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)

# Top gradient bar
add_rect(s, 0, 0, 13.33, 0.06,  C_PURPLE)

# Big background glow circle (decorative)
glow = add_rect(s, 7.8, -1.0, 6.5, 6.5, RGBColor(0x10, 0x0A, 0x28), radius=99999)

# Hackathon tag
tag_pill(s, "🏆  SOLARIS X HACKATHON 2026", 0.65, 1.0, C_PURPLE)

# Title
add_textbox(s, "PersonalOS", 0.6, 1.45, 8.5, 1.3,
            font_size=72, bold=True, color=C_PURPLE_HI)
add_textbox(s, "Agent", 0.6, 2.5, 8.5, 1.2,
            font_size=72, bold=True, color=C_WHITE)

# Tagline
add_textbox(s,
    "Autonomous AI that monitors your inbox, thinks, calls your phone,\n"
    "and acts — with just your voice as the interface.",
    0.6, 3.75, 7.8, 0.85,
    font_size=16, color=C_TEXT2)

# Bullet points
bullets = [
    ("👁️", "3 Autonomous Agents  (Observer · Planner · Executor)"),
    ("📞", "Voice approval via Twilio — fully working"),
    ("📧", "Real Gmail + Google Calendar integration"),
    ("🧠", "ChromaDB memory · MCP protocol · Redis queues"),
]
y = 4.8
for icon, text in bullets:
    add_textbox(s, icon, 0.6,  y, 0.4, 0.32, font_size=13, color=C_WHITE)
    add_textbox(s, text, 1.05, y, 7.5, 0.32, font_size=12, color=C_TEXT2)
    y += 0.36

# Right side — tech stack visual
add_rect(s, 9.3, 1.0, 3.7, 5.8, C_BG2, radius=20000)
add_textbox(s, "TECH STACK", 9.5, 1.15, 3.3, 0.3,
            font_size=8, bold=True, color=C_TEXT3, align=PP_ALIGN.CENTER)
techs = [
    ("Python 3.11 · asyncio",        C_BLUE),
    ("FastAPI · WebSocket",           C_GREEN),
    ("Redis · ChromaDB",              C_ORANGE),
    ("OpenRouter LLM",                C_PURPLE),
    ("MCP (Model Context Protocol)",  C_PURPLE_HI),
    ("Twilio Voice API",              C_RED),
    ("Gmail API · Calendar API",      C_BLUE),
    ("OAuth 2.0 · pydantic",          C_TEXT2),
]
ty = 1.55
for tech, tc in techs:
    add_rect(s, 9.55, ty, 3.2, 0.3, C_BG, radius=15000)
    add_textbox(s, tech, 9.65, ty + 0.04, 3.1, 0.24,
                font_size=9.5, color=tc, align=PP_ALIGN.CENTER)
    ty += 0.38

slide_number_tag(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_textbox(s, "THE PROBLEM", 0.6, 0.22, 6, 0.4,
            font_size=9, bold=True, color=C_TEXT3)
add_textbox(s, "Your inbox owns your time.", 0.6, 0.65, 9, 0.7,
            font_size=36, bold=True, color=C_WHITE)
add_textbox(s,
    "Every day you manually read emails, check your calendar, decide what to reply,\n"
    "open apps, click buttons — for actions that should be automatic.",
    0.6, 1.45, 9, 0.7, font_size=15, color=C_TEXT2)

problems = [
    ("😩", "Urgent email at 2am — you miss it until morning"),
    ("📅", "Meeting request when your calendar is packed — you forget to reply"),
    ("🔁", "Same repetitive replies written over and over"),
    ("🤔", "Deciding priority, urgency, action — mental overhead every single time"),
    ("🔔", "Notification overload — apps, banners, badges competing for attention"),
]
y = 2.4
for icon, text in problems:
    add_rect(s, 0.6, y, 8.5, 0.46, C_BG2, radius=15000)
    add_textbox(s, icon, 0.8,  y + 0.08, 0.4, 0.32, font_size=15)
    add_textbox(s, text, 1.25, y + 0.1, 7.8, 0.3, font_size=12, color=C_TEXT)
    y += 0.56

add_rect(s, 9.4, 1.8, 3.5, 5.3, C_BG2, radius=20000)
add_textbox(s, "The average worker spends", 9.6, 2.1, 3.1, 0.4,
            font_size=11, color=C_TEXT2, align=PP_ALIGN.CENTER)
add_textbox(s, "2.5 hrs", 9.6, 2.55, 3.1, 0.7,
            font_size=44, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
add_textbox(s, "per day on email", 9.6, 3.25, 3.1, 0.35,
            font_size=11, color=C_TEXT2, align=PP_ALIGN.CENTER)
add_rect(s, 9.9, 3.75, 2.5, 0.01, C_BORDER)
add_textbox(s, "28%", 9.6, 3.9, 3.1, 0.7,
            font_size=44, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
add_textbox(s, "of the entire work week", 9.6, 4.6, 3.1, 0.35,
            font_size=11, color=C_TEXT2, align=PP_ALIGN.CENTER)

slide_number_tag(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — THE SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_textbox(s, "THE SOLUTION", 0.6, 0.22, 6, 0.4,
            font_size=9, bold=True, color=C_TEXT3)
add_textbox(s, "An AI that acts, not just advises.", 0.6, 0.65, 10, 0.65,
            font_size=36, bold=True, color=C_WHITE)

add_textbox(s,
    "PersonalOS Agent is a 3-agent autonomous pipeline that monitors your digital life 24/7,\n"
    "reasons about what needs to happen, and calls your phone to confirm — then executes.",
    0.6, 1.38, 12, 0.7, font_size=14, color=C_TEXT2)

agents = [
    ("👁️", "Observer",  "The Eyes",   "Watches Gmail, Calendar,\nFilesystem every 60s via MCP.\nDeduplicates events in Redis.", C_BLUE),
    ("🧠", "Planner",   "The Brain",  "OpenRouter LLM produces\nstructured action plans with\nconfidence 0–100.", C_PURPLE),
    ("⚡", "Executor",  "The Hands",  "Routes by confidence:\nauto-execute, voice approval,\nor silent discard.", C_GREEN),
]
for i, (icon, name, role, desc, col) in enumerate(agents):
    x = 0.6 + i * 4.1
    add_rect(s, x, 2.3, 3.75, 3.0, C_BG2, radius=20000)
    add_rect(s, x, 2.3, 3.75, 0.06, col)
    add_textbox(s, icon, x + 0.15, 2.4,  0.5,  0.5, font_size=26)
    add_textbox(s, name, x + 0.15, 2.95, 3.4,  0.42, font_size=20, bold=True, color=C_WHITE)
    add_textbox(s, role, x + 0.15, 3.38, 3.4,  0.3,  font_size=10, bold=True, color=col)
    add_textbox(s, desc, x + 0.15, 3.72, 3.4,  0.82, font_size=10, color=C_TEXT2)

# Arrows between agents
for ax in [4.35, 8.45]:
    add_textbox(s, "→", ax, 3.5, 0.35, 0.4, font_size=22, color=C_TEXT3, align=PP_ALIGN.CENTER)

# Bottom row — Redis + MCP + ChromaDB
add_rect(s, 0.6, 5.55, 12.1, 0.8, C_BG2, radius=15000)
infra = [
    ("Redis", "Event queues, dedup, activity log", C_RED),
    ("MCP Tool Server", "Gmail · Calendar · Filesystem tools", C_ORANGE),
    ("ChromaDB", "Vector memory · self-improving preferences", C_PURPLE),
    ("Twilio", "Voice calls · speech recognition", C_GREEN),
]
xi = 0.85
for inf_name, inf_desc, inf_col in infra:
    add_textbox(s, inf_name, xi, 5.6, 2.8, 0.26, font_size=10, bold=True, color=inf_col)
    add_textbox(s, inf_desc, xi, 5.88, 2.8, 0.26, font_size=8.5, color=C_TEXT3)
    xi += 3.0

slide_number_tag(s, 3)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_textbox(s, "SYSTEM ARCHITECTURE", 0.6, 0.22, 8, 0.4,
            font_size=9, bold=True, color=C_TEXT3)
add_textbox(s, "How it all fits together", 0.6, 0.62, 8, 0.55,
            font_size=28, bold=True, color=C_WHITE)

# Main pipeline flow
pipeline = [
    (0.55, 1.45, "Gmail / Calendar\nFilesystem", C_BLUE,    "Data\nSources"),
    (2.9,  1.45, "Observer\nAgent",              C_PURPLE,  "Poll &\nNormalize"),
    (5.25, 1.45, "Planner\nAgent",               C_PURPLE,  "LLM\nReasoning"),
    (7.6,  1.45, "Executor\nAgent",              C_GREEN,   "Route &\nExecute"),
    (9.95, 1.45, "MCP Tools\nServer",            C_ORANGE,  "Gmail · Cal\n· Files"),
]
for (bx, by, label, col, sub) in pipeline:
    add_rect(s, bx, by, 2.05, 1.35, C_BG2, radius=20000)
    add_rect(s, bx, by, 2.05, 0.05, col)
    add_textbox(s, label, bx + 0.1, by + 0.15, 1.85, 0.55,
                font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, sub, bx + 0.1, by + 0.72, 1.85, 0.42,
                font_size=9, color=col, align=PP_ALIGN.CENTER)

# Arrows
for ax in [2.62, 4.97, 7.32, 9.67]:
    add_textbox(s, "→", ax, 1.9, 0.3, 0.3, font_size=16, color=C_TEXT3, align=PP_ALIGN.CENTER)

# Redis bus
add_rect(s, 2.9, 3.15, 7.1, 0.6, RGBColor(0x12, 0x0A, 0x28), radius=10000)
add_rect(s, 2.9, 3.15, 7.1, 0.03, C_RED)
add_textbox(s, "Redis Message Bus   ·   events:queue  ·  approvals:pending  ·  dashboard:pending  ·  activity_log",
            3.1, 3.22, 6.8, 0.38, font_size=9, color=C_RED, align=PP_ALIGN.CENTER)

# Lower row
lower = [
    (0.55, 4.0, "ChromaDB\nVector Memory",   C_PURPLE,  "Preferences +\napproval history"),
    (3.05, 4.0, "FastAPI\nDashboard :8080",  C_BLUE,    "WebSocket + REST\nhuman approval UI"),
    (5.55, 4.0, "Twilio\nVoice API",         C_GREEN,   "Outbound calls\nspeech recognition"),
    (8.05, 4.0, "OpenRouter\nLLM",           C_ORANGE,  "Any model:\nGPT-4o · Claude · etc"),
    (10.55,4.0, "MCP / SSE\nProtocol",       C_TEXT2,   "JSON-RPC tools\nstandardized interface"),
]
for (bx, by, label, col, sub) in lower:
    add_rect(s, bx, by, 2.25, 1.3, C_BG2, radius=15000)
    add_textbox(s, label, bx + 0.1, by + 0.12, 2.05, 0.5,
                font_size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(s, sub, bx + 0.1, by + 0.65, 2.05, 0.42,
                font_size=8.5, color=C_TEXT3, align=PP_ALIGN.CENTER)

# Vertical connectors between rows
for cx in [3.95, 6.3, 8.65]:
    add_textbox(s, "↕", cx, 3.8, 0.3, 0.25, font_size=12, color=C_TEXT3, align=PP_ALIGN.CENTER)

# Two-event loop callout
add_rect(s, 0.55, 5.5, 5.0, 0.75, C_BG2, radius=10000)
add_textbox(s, "⚡  Two isolated event loops",
            0.75, 5.55, 4.6, 0.28, font_size=10, bold=True, color=C_ORANGE)
add_textbox(s, "Main loop: agents + dashboard  ·  MCP loop: daemon thread (port 8000)",
            0.75, 5.83, 4.6, 0.28, font_size=8.5, color=C_TEXT3)

add_rect(s, 5.8, 5.5, 5.0, 0.75, C_BG2, radius=10000)
add_textbox(s, "🔄  MCP reconnect per poll cycle",
            6.0, 5.55, 4.6, 0.28, font_size=10, bold=True, color=C_BLUE)
add_textbox(s, "Avoids SSE post_writer drop — fresh session every 60 seconds",
            6.0, 5.83, 4.6, 0.28, font_size=8.5, color=C_TEXT3)

slide_number_tag(s, 4)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — VOICE APPROVAL (THE WOW FEATURE)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
add_rect(s, 0, 0, 13.33, 0.06, C_PURPLE)

add_textbox(s, "🎙️  THE WOW FEATURE", 0.6, 0.18, 10, 0.4,
            font_size=9, bold=True, color=C_PURPLE_HI)
add_textbox(s, "Voice Approval", 0.6, 0.6, 7, 0.65,
            font_size=38, bold=True, color=C_WHITE)
add_textbox(s, "No app. No buttons. Just talk.",
            0.6, 1.28, 7, 0.42, font_size=18, color=C_PURPLE_HI, italic=True)

# Flow steps left column
steps_l = [
    ("📧 Email arrives (70–89% confidence)",           C_BLUE),
    ("Executor stores plan in Redis voice:plan:{id}",  C_TEXT2),
    ("Twilio places outbound call to your phone",      C_GREEN),
    ('You hear: "Requested action: send_email.\n Say yes, no, or modify…"', C_TEXT),
    ("You say: \"reply that I'll submit it by 6pm\"",   C_ORANGE),
]
steps_r = [
    ("Twilio POSTs SpeechResult to /api/twilio/speech/{id}", C_TEXT2),
    ("LLM classifies → MODIFY",                              C_PURPLE_HI),
    ("Voice instruction injected into email body",           C_GREEN),
    ("plan[approved_override] = True → re-queued",           C_TEXT2),
    ("✅ Email sent with your exact words. Goodbye!",         C_GREEN),
]

y = 1.9
for (text, col), (textr, colr) in zip(steps_l, steps_r):
    add_rect(s, 0.55, y, 5.9, 0.5, C_BG2, radius=10000)
    add_rect(s, 0.55, y, 0.04, 0.5, col)
    add_textbox(s, text, 0.72, y + 0.1, 5.6, 0.36, font_size=9.5, color=col)
    if y < 4.0:
        add_textbox(s, "↕", 1.0, y + 0.5, 0.3, 0.28, font_size=11, color=C_TEXT3)

    add_rect(s, 6.85, y, 5.9, 0.5, C_BG2, radius=10000)
    add_rect(s, 6.85, y, 0.04, 0.5, colr)
    add_textbox(s, textr, 7.02, y + 0.1, 5.6, 0.36, font_size=9.5, color=colr)
    if y < 4.0:
        add_textbox(s, "↕", 7.3, y + 0.5, 0.3, 0.28, font_size=11, color=C_TEXT3)

    y += 0.65

# Intent table
add_textbox(s, "INTENT CLASSIFIER — 4 STATES", 0.6, 4.95, 7, 0.32,
            font_size=8.5, bold=True, color=C_TEXT3)
intents = [
    ('"Yes / Go ahead / Sure"',          "APPROVE",  "Executes immediately",               C_GREEN),
    ('"No / Cancel / Reject"',           "REJECT",   "Removed from dashboard",             C_RED),
    ('"Modify … / Reply that … / But…"', "MODIFY",   "Instruction injected into body",     C_ORANGE),
    ("Mumble / Silence / Unclear",       "UNCLEAR",  "Re-prompts you",                     C_TEXT3),
]
xi = 0.55
add_rect(s, xi, 5.28, 12.5, 0.06, C_BORDER)
for speech, intent, action, col in intents:
    add_rect(s, xi, 5.34, 3.1, 0.38, C_BG2)
    add_textbox(s, speech, xi + 0.08, 5.4, 3.0, 0.28, font_size=8.5, color=C_TEXT2, italic=True)
    add_rect(s, xi + 3.12, 5.34, 1.4, 0.38, C_BG2)
    add_textbox(s, intent, xi + 3.16, 5.4, 1.3, 0.28, font_size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_rect(s, xi + 4.55, 5.34, 3.5, 0.38, C_BG2)
    add_textbox(s, action, xi + 4.62, 5.4, 3.4, 0.28, font_size=8.5, color=C_TEXT2)
    add_textbox(s, "→", xi + 3.0, 5.42, 0.2, 0.28, font_size=10, color=C_TEXT3)
    xi += 0.0

# Right side safety callouts
add_rect(s, 8.75, 5.28, 4.3, 1.8, C_BG2, radius=15000)
add_textbox(s, "Safety features", 8.9, 5.35, 4.0, 0.3, font_size=10, bold=True, color=C_PURPLE_HI)
safeties = [
    "• 10s speechTimeout — plenty of time to speak",
    "• Word-boundary regex: 'no' ≠ 'know' / 'notify'",
    "• MODIFY checked before REJECT (no false triggers)",
    "• LLM fails → keyword fallback activates",
    "• Crash in handler → graceful TwiML error reply",
]
sy = 5.68
for safe in safeties:
    add_textbox(s, safe, 8.9, sy, 4.05, 0.25, font_size=8.5, color=C_TEXT2)
    sy += 0.25

slide_number_tag(s, 5)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — VOICE DEMO SCREENSHOT (SIMULATED)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s, C_GREEN)

add_textbox(s, "LIVE DEMO — VOICE FLOW", 0.6, 0.22, 8, 0.4,
            font_size=9, bold=True, color=C_GREEN)
add_textbox(s, "What actually happened", 0.6, 0.62, 8, 0.55,
            font_size=28, bold=True, color=C_WHITE)
add_textbox(s, "Confirmed working — email received at vikas935314@gmail.com with voice instruction in body",
            0.6, 1.18, 12, 0.38, font_size=12, color=C_GREEN, italic=True)

# Terminal window simulation
add_rect(s, 0.55, 1.7, 7.3, 4.8, RGBColor(0x03, 0x05, 0x0E), radius=15000)
add_rect(s, 0.55, 1.7, 7.3, 0.32, RGBColor(0x1A, 0x1F, 0x35), radius=15000)
for cx, cc in [(0.82, C_RED), (1.1, C_ORANGE), (1.38, C_GREEN)]:
    add_rect(s, cx, 1.78, 0.12, 0.12, cc, radius=99999)

terminal_lines = [
    ("[Executor] Ready — voice approval: ON",                  C_GREEN),
    ("[Executor] DASHBOARD APPROVAL REQUIRED",                 C_ORANGE),
    ("  Action    : send_email",                              C_TEXT2),
    ("  Confidence: 78%",                                     C_TEXT2),
    ("[Executor] VOICE CALL dispatched — Plan: a3f9b2...",    C_PURPLE_HI),
    ("",                                                       C_TEXT3),
    ("[VoiceApproval] Speech received — plan=a3f9b2",         C_BLUE),
    ("  speech='reply that I will submit it by 6pm today'",   C_TEXT),
    ("",                                                       C_TEXT3),
    ("[VoiceApproval] LLM classified → MODIFY",               C_ORANGE),
    ("  instruction: 'I will submit it by 6pm today'",        C_TEXT2),
    ("  → injecting into action_args.body",                   C_TEXT2),
    ("",                                                       C_TEXT3),
    ("[Executor] AUTO-EXECUTE: send_email",                    C_GREEN),
    ("  Result: {'message_id': '18f...', 'status': 'sent'}",  C_GREEN),
    ("[VoiceApproval] voice_approved plan=a3f9b2",             C_GREEN),
]
ty = 2.15
for line_text, line_col in terminal_lines:
    if ty > 6.2:
        break
    add_textbox(s, line_text, 0.72, ty, 7.0, 0.24,
                font_size=8.5, color=line_col)
    ty += 0.24

# Email preview
add_rect(s, 8.1, 1.7, 5.0, 4.8, C_BG2, radius=15000)
add_rect(s, 8.1, 1.7, 5.0, 0.04, C_GREEN)
add_textbox(s, "📧  EMAIL SENT", 8.3, 1.8, 4.6, 0.3,
            font_size=9.5, bold=True, color=C_GREEN)
add_rect(s, 8.25, 2.15, 4.65, 0.01, C_BORDER)

fields = [
    ("TO",      "vikas935314@gmail.com"),
    ("SUBJECT", "PersonalOS Agent — Voice Approved"),
    ("STATUS",  "✓ Sent via Gmail API"),
]
fy = 2.25
for label, val in fields:
    add_textbox(s, label, 8.3, fy, 0.9, 0.26,
                font_size=7.5, bold=True, color=C_TEXT3)
    add_textbox(s, val, 9.25, fy, 3.7, 0.26,
                font_size=8.5, color=C_TEXT)
    fy += 0.3

add_rect(s, 8.25, 3.1, 4.65, 0.01, C_BORDER)
add_textbox(s, "EMAIL BODY:", 8.3, 3.2, 4.5, 0.26,
            font_size=7.5, bold=True, color=C_TEXT3)

email_body = (
    "Hi,\n\n"
    "I will submit it by 6pm today,\n"
    "please give me some more time\n\n"
    "---\n"
    "Sent via PersonalOS Agent\n"
    "(voice approved)\n"
    "Powered by SOLARIS X"
)
add_textbox(s, email_body, 8.3, 3.5, 4.65, 2.75,
            font_size=9.5, color=C_TEXT2, italic=True)

add_rect(s, 8.1, 6.2, 5.0, 0.22, RGBColor(0x06, 0x2E, 0x18), radius=10000)
add_textbox(s, "⚡ Delivered in < 3 seconds after spoken approval",
            8.2, 6.22, 4.8, 0.18, font_size=8, color=C_GREEN, align=PP_ALIGN.CENTER)

slide_number_tag(s, 6)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — MEETING INTELLIGENCE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s, C_BLUE)

add_textbox(s, "NEW FEATURE", 0.6, 0.22, 5, 0.4,
            font_size=9, bold=True, color=C_BLUE)
add_textbox(s, "Smart Meeting Intelligence", 0.6, 0.62, 10, 0.62,
            font_size=32, bold=True, color=C_WHITE)
add_textbox(s,
    "Reads your calendar, detects conflicts, and auto-drafts the right reply — for every scenario.",
    0.6, 1.3, 12, 0.4, font_size=14, color=C_TEXT2)

scenarios = [
    ("1", "Direct conflict",      '"Can we meet tomorrow 3pm?"',  "3pm slot blocked",          '"Already have something at 3pm. 5pm work?"',          C_RED),
    ("2", "Full day packed",      '"Free anytime Thursday?"',     "All-day back-to-back",       '"Thursday is fully packed. Friday 10am?"',            C_RED),
    ("3", "Partial overlap",      '"1hr call at 2pm Friday"',     "2:30pm conflict",            '"Conflict at 2:30. 1pm instead, or Monday?"',         C_ORANGE),
    ("4", "No time specified",    '"Catch up this week"',         "Wed PM + Fri AM free",       '"Free Wed after 3pm or Fri before noon"',             C_GREEN),
    ("5", "Next day also busy",   '"Reschedule to tomorrow?"',    "Tomorrow blocked too",       '"Tomorrow also packed. Thursday 10am?"',              C_RED),
    ("6", "Recurring conflict",   '"Every Monday 2pm?"',          "Weekly Sync on Mon 2pm",     '"Mondays blocked. Tue or Wed same time?"',            C_ORANGE),
    ("7", "Calendar is FREE",     '"Meet tomorrow at 11am?"',     "11am–12pm free",             '"Confirmed + calendar invite created!"',              C_GREEN),
]

y = 1.85
for i, (num, name, trigger, cal, reply, col) in enumerate(scenarios):
    row = i % 4
    col_x = 0.55 + (i // 4) * 6.65
    if i == 4:
        y = 1.85
    if i < 4:
        y = 1.85 + i * 1.3
    else:
        y = 1.85 + (i - 4) * 1.3

    # clamp to 3 per column, 4 on left
    if i < 4:
        bx = 0.55
        by = 1.85 + i * 1.35
    else:
        bx = 6.65
        by = 1.85 + (i - 4) * 1.35

    add_rect(s, bx, by, 5.85, 1.18, C_BG2, radius=12000)
    add_rect(s, bx, by, 5.85, 0.04, col)
    # Number badge
    add_rect(s, bx + 0.1, by + 0.1, 0.28, 0.28, col, radius=99999)
    add_textbox(s, num, bx + 0.1, by + 0.11, 0.28, 0.24,
                font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, name, bx + 0.48, by + 0.1, 5.0, 0.26,
                font_size=10.5, bold=True, color=C_WHITE)
    add_textbox(s, f"📧 {trigger}", bx + 0.12, by + 0.42, 5.5, 0.24,
                font_size=8.5, color=C_TEXT2, italic=True)
    add_textbox(s, f"📅 {cal}", bx + 0.12, by + 0.66, 5.5, 0.24,
                font_size=8.5, color=C_TEXT3)
    add_textbox(s, f"✉️ {reply}", bx + 0.12, by + 0.9, 5.5, 0.24,
                font_size=8.5, color=col)

slide_number_tag(s, 7)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — LLM + DECISION SCORING
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s, C_PURPLE)

add_textbox(s, "INTELLIGENCE LAYER", 0.6, 0.22, 8, 0.38,
            font_size=9, bold=True, color=C_TEXT3)
add_textbox(s, "LLM Planning + Decision Scoring", 0.6, 0.62, 10, 0.6,
            font_size=30, bold=True, color=C_WHITE)

# Scoring formula card
add_rect(s, 0.55, 1.4, 6.2, 1.55, C_BG2, radius=15000)
add_rect(s, 0.55, 1.4, 6.2, 0.05, C_PURPLE)
add_textbox(s, "DECISION SCORING FORMULA", 0.75, 1.5, 5.8, 0.3,
            font_size=8.5, bold=True, color=C_TEXT3)
add_textbox(s, "adjusted = base × urgency_mult × history_mult",
            0.75, 1.88, 5.8, 0.38, font_size=14, bold=True, color=C_PURPLE_HI)
add_textbox(s, "urgency_mult : 1.0 – 1.3  (+10% per urgency keyword, max +30%)",
            0.75, 2.3, 5.8, 0.26, font_size=9.5, color=C_TEXT2)
add_textbox(s, "history_mult : 0.9 – 1.1  (ChromaDB historical approval rate)",
            0.75, 2.56, 5.8, 0.26, font_size=9.5, color=C_TEXT2)

# Routing card
add_rect(s, 0.55, 3.1, 6.2, 1.65, C_BG2, radius=15000)
add_rect(s, 0.55, 3.1, 6.2, 0.05, C_GREEN)
add_textbox(s, "CONFIDENCE ROUTING", 0.75, 3.2, 5.8, 0.3,
            font_size=8.5, bold=True, color=C_TEXT3)
routing = [
    ("> 90%",    "Auto-execute + Twilio notification",          C_GREEN),
    ("70 – 89%", "Voice approval call + dashboard pending",     C_ORANGE),
    ("< 70%",    "Silent discard (logged, learned from)",       C_RED),
]
ry = 3.55
for conf, action, col in routing:
    add_rect(s, 0.75, ry, 1.3, 0.34, col, radius=15000)
    add_textbox(s, conf, 0.75, ry + 0.04, 1.3, 0.28,
                font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, action, 2.18, ry + 0.05, 4.4, 0.28,
                font_size=10, color=C_TEXT2)
    ry += 0.42

# LLM prompt schema
add_rect(s, 7.0, 1.4, 5.9, 5.35, RGBColor(0x03, 0x05, 0x0E), radius=15000)
add_textbox(s, "JSON PLAN SCHEMA  (LLM must follow exactly)", 7.2, 1.52, 5.5, 0.28,
            font_size=8.5, bold=True, color=C_TEXT3)
schema_text = '''{
  "action": "send_email",
  "confidence": 78,
  "priority": "high",
  "reason": "Professor asking for urgent reply",
  "requires_approval": true,
  "alternatives": [
    {"action": "no_action", "confidence": 20,
     "reason": "Wait and monitor"},
    {"action": "read_emails", "confidence": 45,
     "reason": "Fetch full thread first"}
  ],
  "explanation": "High urgency keywords detected.
    Historical approval rate: 0.85.
    Meeting context injected from calendar.",
  "action_args": {
    "to": "professor@university.edu",
    "subject": "Re: Assignment Deadline",
    "body": "Hi Professor, ..."
  }
}'''
add_textbox(s, schema_text, 7.15, 1.88, 5.6, 4.55,
            font_size=8, color=C_TEXT2)

# ChromaDB memory callout
add_rect(s, 0.55, 4.9, 6.2, 1.85, C_BG2, radius=15000)
add_rect(s, 0.55, 4.9, 6.2, 0.05, C_PURPLE)
add_textbox(s, "CHROMADB VECTOR MEMORY", 0.75, 5.0, 5.8, 0.3,
            font_size=8.5, bold=True, color=C_TEXT3)
add_textbox(s, "Self-improving over time", 0.75, 5.32, 5.8, 0.32,
            font_size=13, bold=True, color=C_PURPLE_HI)
mem_items = [
    "Stores user preference statements as embeddings",
    "Semantically queried per event — personalises every plan",
    "Records approval/rejection outcomes per action type",
    "Adjusts history_mult: good history → higher confidence",
]
my = 5.68
for item in mem_items:
    add_textbox(s, f"• {item}", 0.75, my, 5.8, 0.25, font_size=9, color=C_TEXT2)
    my += 0.25

slide_number_tag(s, 8)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s, C_BLUE)

add_textbox(s, "HUMAN-IN-THE-LOOP", 0.6, 0.22, 8, 0.38,
            font_size=9, bold=True, color=C_BLUE)
add_textbox(s, "Approval Dashboard", 0.6, 0.62, 8, 0.58,
            font_size=32, bold=True, color=C_WHITE)
add_textbox(s, "FastAPI + WebSocket · Light/Dark theme · Real-time updates",
            0.6, 1.25, 10, 0.36, font_size=13, color=C_TEXT2)

# 4 tabs
tabs = [
    ("📬", "Email Intelligence",
     ["All email plans with priority badges",
      "Search by subject, sender, action",
      "Filter: high / medium / pending / auto",
      "Approve / Reject per item or Bulk All",
      "Export CSV for records"]),
    ("📊", "Analytics",
     ["Confidence distribution histogram",
      "Priority breakdown bar chart",
      "Response breakdown (auto/pending/rejected)",
      "Auto-execute rate indicator",
      "Queue depth monitor"]),
    ("⚡", "Inject Event",
     ["Push synthetic email / calendar event",
      "No Google OAuth required",
      "Mark as urgent — boosts confidence",
      "Perfect for live hackathon demos",
      "Flows through full pipeline in real time"]),
    ("🧠", "Preferences",
     ["View all ChromaDB user preferences",
      "Add natural-language statements",
      "e.g. 'Always reply within 2 hours'",
      "Semantically queried per planning event",
      "Personalises every LLM decision"]),
]
tx = 0.55
for icon, title, points in tabs:
    add_rect(s, tx, 1.75, 3.0, 4.85, C_BG2, radius=15000)
    add_rect(s, tx, 1.75, 3.0, 0.05, C_BLUE)
    add_textbox(s, icon, tx + 0.12, 1.85, 0.5, 0.45, font_size=22)
    add_textbox(s, title, tx + 0.12, 2.38, 2.7, 0.36,
                font_size=11.5, bold=True, color=C_WHITE)
    py = 2.82
    for pt in points:
        add_textbox(s, f"• {pt}", tx + 0.12, py, 2.75, 0.3,
                    font_size=9, color=C_TEXT2)
        py += 0.3
    tx += 3.2

# Bottom features
bottom = [
    ("⚡ WebSocket", "Real-time push — no polling needed"),
    ("⌨️ Keyboard nav", "A approve · R reject · ↑↓ navigate"),
    ("📞 Test Call btn", "Trigger Twilio demo from header"),
    ("🌙/☀️ Theme toggle", "Persisted in localStorage"),
]
bx = 0.55
for bname, bdesc in bottom:
    add_rect(s, bx, 6.75, 3.0, 0.55, C_BG2, radius=10000)
    add_textbox(s, bname, bx + 0.12, 6.8, 2.7, 0.24,
                font_size=9.5, bold=True, color=C_PURPLE_HI)
    add_textbox(s, bdesc, bx + 0.12, 7.02, 2.7, 0.22,
                font_size=8.5, color=C_TEXT3)
    bx += 3.2

slide_number_tag(s, 9)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — LIVE DEMO FLOW
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s, C_ORANGE)

add_textbox(s, "LIVE DEMO", 0.6, 0.22, 8, 0.38,
            font_size=9, bold=True, color=C_ORANGE)
add_textbox(s, "Watch it work in 90 seconds", 0.6, 0.62, 10, 0.58,
            font_size=32, bold=True, color=C_WHITE)

steps = [
    ("1", "python main.py",
     "All 3 agents + MCP server + dashboard start in ~3 seconds.\nHealth strip shows all services green.",
     C_GREEN),
    ("2", "Open localhost:8080",
     "Dashboard loads. Health strip: Redis ✓, LLM ✓, Google ✓, Twilio ✓\nToggle dark/light theme to show polish.",
     C_BLUE),
    ("3", "Click 'Inject Event'",
     "From: professor@university.edu\nSubject: Can we meet tomorrow at 3pm? — URGENT\nCheck 'Mark as urgent' → click Inject Event →",
     C_ORANGE),
    ("4", "Watch the pipeline fire",
     "Activity feed: Observer detected → Planner reasoning → calendar checked\nEmail Intelligence tab updates live via WebSocket",
     C_PURPLE),
    ("5", "Phone rings",
     "Answer it. Hear the AI speak the action.\nSay: 'Modify, reply that I'm free at 5pm instead'\nHear: 'Got it! Goodbye!'",
     C_GREEN),
    ("6", "Show the result",
     "Dashboard: 'VOICE MODIFIED+APPROVED' badge\nCheck Gmail — email sent with your exact spoken words in body.\nAnalytics tab shows confidence scores.",
     C_ORANGE),
]
y = 1.42
for i, (num, title, desc, col) in enumerate(steps):
    bx = 0.55 if i < 3 else 6.8
    by = 1.42 + (i % 3) * 1.98
    add_rect(s, bx, by, 5.9, 1.78, C_BG2, radius=15000)
    add_rect(s, bx, by, 0.045, 1.78, col)
    # Step number
    add_rect(s, bx + 0.18, by + 0.18, 0.42, 0.42, col, radius=99999)
    add_textbox(s, num, bx + 0.18, by + 0.19, 0.42, 0.38,
                font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, title, bx + 0.74, by + 0.18, 4.9, 0.38,
                font_size=13, bold=True, color=C_WHITE)
    add_textbox(s, desc, bx + 0.18, by + 0.68, 5.55, 0.9,
                font_size=9.5, color=C_TEXT2)

slide_number_tag(s, 10)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — MCP PROTOCOL
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s, C_ORANGE)

add_textbox(s, "TECHNICAL DEPTH", 0.6, 0.22, 8, 0.38,
            font_size=9, bold=True, color=C_TEXT3)
add_textbox(s, "MCP — Model Context Protocol", 0.6, 0.62, 10, 0.58,
            font_size=30, bold=True, color=C_WHITE)
add_textbox(s,
    "The same protocol Anthropic built into Claude. Agents call tools over SSE JSON-RPC — "
    "fully decoupled from any specific API implementation.",
    0.6, 1.28, 12.1, 0.5, font_size=13, color=C_TEXT2)

tools = [
    ("read_emails",   "max_results, query",                           "List of Gmail messages",       C_BLUE),
    ("send_email",    "to, subject, body",                            "{message_id, status}",         C_BLUE),
    ("read_calendar", "days_ahead",                                   "List of Calendar events",      C_GREEN),
    ("create_event",  "summary, start_datetime, end_datetime, ...",   "{event_id, html_link}",        C_GREEN),
    ("list_files",    "directory",                                    "List of file metadata",        C_ORANGE),
    ("move_file",     "source, destination",                          "{source, dest, success}",      C_ORANGE),
]
add_textbox(s, "TOOL", 0.55, 2.0, 2.8, 0.28, font_size=8, bold=True, color=C_TEXT3)
add_textbox(s, "ARGUMENTS", 3.5, 2.0, 4.2, 0.28, font_size=8, bold=True, color=C_TEXT3)
add_textbox(s, "RETURNS", 7.85, 2.0, 3.5, 0.28, font_size=8, bold=True, color=C_TEXT3)
add_rect(s, 0.55, 2.28, 11.5, 0.02, C_BORDER)

ty = 2.35
for tool_name, args, returns, col in tools:
    add_rect(s, 0.55, ty, 11.5, 0.42, C_BG2)
    add_textbox(s, tool_name, 0.65, ty + 0.08, 2.75, 0.28,
                font_size=10, bold=True, color=col)
    add_textbox(s, args, 3.5, ty + 0.08, 4.2, 0.28, font_size=9, color=C_TEXT2)
    add_textbox(s, returns, 7.85, ty + 0.08, 3.5, 0.28, font_size=9, color=C_TEXT3)
    ty += 0.44

# Right: architecture callouts
add_rect(s, 0.55, 5.2, 11.5, 2.0, C_BG2, radius=15000)
add_textbox(s, "WHY MCP MATTERS FOR THIS PROJECT", 0.75, 5.3, 11.0, 0.3,
            font_size=8.5, bold=True, color=C_TEXT3)
reasons = [
    ("🔌 Swappable tools",     "Swap Gmail → Outlook, or Google Cal → iCloud Cal without touching agent code"),
    ("🔒 Security boundary",   "Agents never hold credentials — only MCP server does"),
    ("🔄 Retry isolation",     "Tool retries happen at BaseAgent layer — agents don't care about flaky network"),
    ("📊 Auditable",           "Every tool call flows through one server — easy to log, monitor, and rate-limit"),
]
ry = 5.65
for rname, rdesc in reasons:
    add_textbox(s, rname, 0.75, ry, 2.2, 0.28, font_size=9.5, bold=True, color=C_ORANGE)
    add_textbox(s, rdesc, 3.1, ry, 9.0, 0.28, font_size=9, color=C_TEXT2)
    ry += 0.35

slide_number_tag(s, 11)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — FEATURE GRID
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_textbox(s, "WHAT WE BUILT", 0.6, 0.22, 8, 0.38,
            font_size=9, bold=True, color=C_TEXT3)
add_textbox(s, "18 production-grade features", 0.6, 0.62, 10, 0.58,
            font_size=30, bold=True, color=C_WHITE)

features = [
    ("📧", "Real Gmail monitoring",        "Unread poll, dedup across restarts",   C_BLUE),
    ("📅", "Google Calendar",              "Read + create events, OAuth shared",   C_BLUE),
    ("🧠", "LLM action planning",          "Any OpenRouter model, strict JSON",    C_PURPLE),
    ("📊", "Confidence routing",           ">90 auto, 70–89 voice, <70 discard",  C_GREEN),
    ("📞", "Voice approval",               "Full Twilio pipeline, 10s timeout",    C_GREEN),
    ("✏️", "Voice modification",           "Spoken words injected into email",     C_GREEN),
    ("🗓️", "Meeting intelligence",         "7 conflict scenarios, free-slot AI",   C_BLUE),
    ("🖥️", "Approval dashboard",           "FastAPI, WebSocket, 4 tabs",           C_PURPLE),
    ("⚡", "Real-time WebSocket",          "Instant push, 20s keep-alive ping",    C_ORANGE),
    ("💾", "ChromaDB memory",              "Vector prefs + approval history",      C_PURPLE),
    ("📈", "Decision scoring",             "base × urgency × history formula",     C_ORANGE),
    ("🔧", "MCP protocol",                 "6 tools, SSE JSON-RPC, swappable",     C_ORANGE),
    ("🗂️", "Filesystem assistant",         "Sandboxed list + move, path guard",    C_TEXT2),
    ("💉", "Inject Event demo mode",       "No OAuth needed for live demos",       C_GREEN),
    ("🧪", "33 automated tests",           "18 unit + 15 integration, all pass",   C_RED),
    ("🌙", "Light/Dark theme",             "CSS vars, persisted localStorage",     C_TEXT2),
    ("📝", "Structlog + JSON",             "Colorized console + file logging",     C_TEXT2),
    ("🛡️", "Graceful shutdown",            "Clean MCP SSE close, no leak",         C_TEXT2),
]
cols = 3
fw = 4.25
fh = 0.72
margin_x = 0.5
margin_y = 1.4
for i, (icon, name, desc, col) in enumerate(features):
    row = i // cols
    column = i % cols
    bx = margin_x + column * (fw + 0.06)
    by = margin_y + row * (fh + 0.06)
    add_rect(s, bx, by, fw, fh, C_BG2, radius=12000)
    add_rect(s, bx, by, fw, 0.04, col)
    add_textbox(s, icon, bx + 0.1, by + 0.09, 0.38, 0.35, font_size=14)
    add_textbox(s, name, bx + 0.52, by + 0.1, fw - 0.65, 0.3,
                font_size=10, bold=True, color=C_WHITE)
    add_textbox(s, desc, bx + 0.52, by + 0.4, fw - 0.65, 0.26,
                font_size=8.5, color=C_TEXT3)

slide_number_tag(s, 12)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — WHY WE WIN
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s, C_PURPLE)

add_textbox(s, "WHY WE WIN", 0.6, 0.22, 8, 0.38,
            font_size=9, bold=True, color=C_PURPLE_HI)
add_textbox(s, "What makes this different", 0.6, 0.62, 10, 0.58,
            font_size=32, bold=True, color=C_WHITE)

reasons = [
    ("📞", "Voice is REAL",
     "Not a mockup. Not a simulation.\nThe phone rings. You answer. You speak.\nThe email arrives. Confirmed on a real Gmail inbox.",
     C_GREEN),
    ("🏗️", "Production quality",
     "33 passing tests. Structured logging.\nRetry logic with exponential backoff.\nPre-flight checker. Graceful shutdown.",
     C_BLUE),
    ("🔗", "Real integrations",
     "Live Gmail API. Live Google Calendar.\nLive Twilio. Authenticated OAuth2.\nNot mock data — real emails flowing.",
     C_ORANGE),
    ("🧠", "Self-improving AI",
     "ChromaDB records every approval.\nFuture confidence adjusts from history.\nGets smarter with every interaction.",
     C_PURPLE),
    ("🔌", "Industry-standard MCP",
     "Same protocol Anthropic built into Claude.\nAgents are fully decoupled from APIs.\nSwap any tool without touching agents.",
     C_PURPLE_HI),
    ("💡", "Novel UX model",
     "AI monitors → reasons → calls you → listens → acts.\nNo app to open. No button to press.\nYour phone IS the interface.",
     C_RED),
]
for i, (icon, title, desc, col) in enumerate(reasons):
    row = i // 3
    col_n = i % 3
    bx = 0.55 + col_n * 4.24
    by = 1.45 + row * 2.5
    add_rect(s, bx, by, 3.95, 2.22, C_BG2, radius=15000)
    add_rect(s, bx, by, 3.95, 0.05, col)
    add_textbox(s, icon, bx + 0.15, by + 0.15, 0.5, 0.4, font_size=22)
    add_textbox(s, title, bx + 0.72, by + 0.16, 3.05, 0.36,
                font_size=13, bold=True, color=col)
    add_textbox(s, desc, bx + 0.15, by + 0.62, 3.65, 1.35,
                font_size=9.5, color=C_TEXT2)

slide_number_tag(s, 13)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
add_rect(s, 0, 0, 13.33, 0.06, C_PURPLE)
add_rect(s, 0, 7.44, 13.33, 0.06, C_PURPLE)

add_textbox(s, "SOLARIS X HACKATHON 2026", 0, 0.9, 13.33, 0.5,
            font_size=10, bold=True, color=C_TEXT3, align=PP_ALIGN.CENTER)
add_textbox(s, "PersonalOS Agent", 0, 1.55, 13.33, 0.95,
            font_size=56, bold=True, color=C_PURPLE_HI, align=PP_ALIGN.CENTER)

add_textbox(s,
    '"An autonomous AI that monitors your inbox, reasons about what needs to happen,\n'
    'calls your phone to ask — and when you say  \'reply that I\'ll submit it by evening\' —\n'
    'it sends the email with those exact words."',
    0.8, 2.65, 11.7, 1.1, font_size=15.5, color=C_TEXT2, align=PP_ALIGN.CENTER, italic=True)

add_textbox(s, "No app.  No buttons.  Just talk.",
            0, 3.9, 13.33, 0.55, font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Stats row
stats = [
    ("3",  "Autonomous\nAgents",     C_PURPLE_HI),
    ("6",  "MCP\nTools",             C_BLUE),
    ("7",  "Meeting\nScenarios",     C_GREEN),
    ("4",  "Voice\nIntents",         C_ORANGE),
    ("33", "Tests\nPassing",         C_RED),
    ("∞",  "Possibilities",          C_PURPLE_HI),
]
sx = 0.75
for snum, slabel, scol in stats:
    add_rect(s, sx, 4.7, 1.85, 1.25, C_BG2, radius=15000)
    add_textbox(s, snum, sx, 4.8, 1.85, 0.6,
                font_size=34, bold=True, color=scol, align=PP_ALIGN.CENTER)
    add_textbox(s, slabel, sx, 5.42, 1.85, 0.45,
                font_size=8.5, bold=True, color=C_TEXT3, align=PP_ALIGN.CENTER)
    sx += 1.97

add_textbox(s, "Built with Python · FastAPI · Redis · ChromaDB · MCP · Twilio · Gmail API · Google Calendar API · OpenRouter",
            0, 6.25, 13.33, 0.35, font_size=9.5, color=C_TEXT3, align=PP_ALIGN.CENTER)
add_textbox(s, "🏆  Thank you — let's change how people manage their digital lives.",
            0, 6.72, 13.33, 0.4, font_size=12, bold=True, color=C_PURPLE_HI, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "docs", "PersonalOS_Agent_Presentation.pptx")
prs.save(out)
print(f"OK  Saved: {out}")
print(f"    Slides: {len(prs.slides)}")
